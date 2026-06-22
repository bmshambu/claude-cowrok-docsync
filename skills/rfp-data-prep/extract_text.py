"""
extract_text.py
---------------
Extracts text + page/paragraph references from PDF, DOCX, and PPTX files
in rfp_docs/ and writes:
  - extracted_text/<docname>.txt      (full plain text)
  - chunks/<docname>_chunks.json      (chunked with page/para refs for citation)

Usage:
    python scripts/extract_text.py
"""

import os, json, re
from pathlib import Path

# ── Paths (relative to Smart_RAG root) ──────────────────────────────────────
ROOT        = Path(__file__).parent.parent.parent
RFP_DIR     = ROOT / "rfp_docs"
TEXT_DIR    = ROOT / "extracted_text"
CHUNKS_DIR  = ROOT / "chunks"

TEXT_DIR.mkdir(exist_ok=True)
CHUNKS_DIR.mkdir(exist_ok=True)

CHUNK_SIZE  = 400   # words per chunk
OVERLAP     = 50    # word overlap between chunks


# ── Helpers ──────────────────────────────────────────────────────────────────

def clean(text: str) -> str:
    """Remove excessive whitespace."""
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def make_chunks(paragraphs: list[dict], doc_id: str, filename: str) -> list[dict]:
    """
    paragraphs: list of {"text": str, "page": int|None, "section": str}
    Returns chunks with word-level overlap and source refs for citation.
    """
    chunks = []
    chunk_id = 0
    buffer_words = []
    buffer_meta  = []   # (word_index, page, section)

    for para in paragraphs:
        words = para["text"].split()
        for w in words:
            buffer_meta.append((len(buffer_words), para["page"], para["section"]))
            buffer_words.append(w)

        # flush when buffer reaches chunk size
        while len(buffer_words) >= CHUNK_SIZE:
            chunk_words = buffer_words[:CHUNK_SIZE]
            meta_slice  = buffer_meta[:CHUNK_SIZE]

            # derive page + section from first token in chunk
            first_page    = meta_slice[0][1]
            first_section = meta_slice[0][2]
            last_page     = meta_slice[-1][1]

            chunks.append({
                "chunk_id"   : f"{doc_id}_chunk_{chunk_id:04d}",
                "doc_id"     : doc_id,
                "filename"   : filename,
                "page_start" : first_page,
                "page_end"   : last_page,
                "section"    : first_section,
                "text"       : " ".join(chunk_words)
            })
            chunk_id += 1
            # slide window with overlap
            buffer_words = buffer_words[CHUNK_SIZE - OVERLAP:]
            buffer_meta  = buffer_meta[CHUNK_SIZE - OVERLAP:]

    # flush remainder
    if buffer_words:
        meta_slice = buffer_meta
        chunks.append({
            "chunk_id"   : f"{doc_id}_chunk_{chunk_id:04d}",
            "doc_id"     : doc_id,
            "filename"   : filename,
            "page_start" : meta_slice[0][1]  if meta_slice else None,
            "page_end"   : meta_slice[-1][1] if meta_slice else None,
            "section"    : meta_slice[0][2]  if meta_slice else "",
            "text"       : " ".join(buffer_words)
        })

    return chunks


# ── Format extractors ─────────────────────────────────────────────────────────

def extract_docx(path: Path) -> list[dict]:
    """Returns list of {text, page, section}. DOCX has no page concept — use para index."""
    from docx import Document
    doc = Document(path)
    paragraphs = []
    current_section = "Document Start"
    page_estimate = 1
    word_count = 0

    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if not text:
            continue

        # detect section headings
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            current_section = text

        # rough page estimate: ~300 words per page
        word_count += len(text.split())
        page_estimate = max(1, word_count // 300 + 1)

        paragraphs.append({
            "text"    : text,
            "page"    : page_estimate,
            "section" : current_section
        })

    return paragraphs


def extract_pdf(path: Path) -> list[dict]:
    import fitz  # PyMuPDF
    doc = fitz.open(path)
    paragraphs = []
    current_section = "Document Start"

    for page_num, page in enumerate(doc, start=1):
        blocks = page.get_text("blocks")
        for block in blocks:
            text = block[4].strip()
            if not text:
                continue
            # heuristic: short ALL-CAPS or bold-ish lines = section heading
            if len(text) < 80 and text.isupper():
                current_section = text
            paragraphs.append({
                "text"    : clean(text),
                "page"    : page_num,
                "section" : current_section
            })

    return paragraphs


def extract_pptx(path: Path) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(path)
    paragraphs = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        current_section = f"Slide {slide_num}"
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                paragraphs.append({
                    "text"    : text,
                    "page"    : slide_num,   # slide number as "page"
                    "section" : current_section
                })

    return paragraphs


# ── Main ─────────────────────────────────────────────────────────────────────

def process_file(path: Path):
    ext    = path.suffix.lower()
    doc_id = path.stem.replace(" ", "_")

    print(f"  Processing: {path.name}")

    if ext == ".docx":
        paragraphs = extract_docx(path)
    elif ext == ".pdf":
        paragraphs = extract_pdf(path)
    elif ext in (".pptx", ".ppt"):
        paragraphs = extract_pptx(path)
    else:
        print(f"    ⚠ Skipping unsupported format: {ext}")
        return

    # Write full plain text
    full_text = "\n\n".join(p["text"] for p in paragraphs)
    (TEXT_DIR / f"{doc_id}.txt").write_text(full_text, encoding="utf-8")

    # Write chunks JSON
    chunks = make_chunks(paragraphs, doc_id, path.name)
    (CHUNKS_DIR / f"{doc_id}_chunks.json").write_text(
        json.dumps(chunks, indent=2, ensure_ascii=False), encoding="utf-8"
    )

    print(f"    ✓ {len(paragraphs)} paragraphs → {len(chunks)} chunks")
    return {"doc_id": doc_id, "filename": path.name, "chunks": len(chunks), "paragraphs": len(paragraphs)}


def main():
    files = list(RFP_DIR.glob("*"))
    supported = [f for f in files if f.suffix.lower() in (".pdf", ".docx", ".pptx", ".ppt")]

    print(f"\n📁 Found {len(supported)} supported file(s) in rfp_docs/\n")

    results = []
    for f in supported:
        r = process_file(f)
        if r:
            results.append(r)

    print(f"\n✅ Extraction complete.")
    print(f"   Files processed : {len(results)}")
    print(f"   extracted_text/ : {TEXT_DIR}")
    print(f"   chunks/         : {CHUNKS_DIR}")
    print("\nNext step: run entity extraction via Claude Cowork on each file in extracted_text/")


if __name__ == "__main__":
    main()
